"""
Build the project pitch deck (PowerPoint) from the real results.

Figures are the actual plots in results/. Run:
    .venv/bin/python paper/build_deck.py
        → paper/RL_Routing_Pitch.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE = os.path.dirname(__file__)
ROOT = os.path.dirname(HERE)
RESULTS = os.path.join(ROOT, "results")
OUT = os.path.join(HERE, "RL_Routing_Pitch.pptx")

NAVY = RGBColor(0x12, 0x2A, 0x4A)
ACCENT = RGBColor(0x6A, 0x1B, 0x9A)
GREY = RGBColor(0x44, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def textbox(s, left, top, width, height):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def title_bar(s, text, sub=None):
    tf = textbox(s, 0.55, 0.35, 12.2, 1.1)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(15); r2.font.color.rgb = ACCENT; r2.font.italic = True
    # accent underline
    ln = s.shapes.add_shape(1, Inches(0.6), Inches(1.45), Inches(12.1), Pt(2.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()


def bullets(s, items, left=0.7, top=1.8, width=7.0, height=5.0, size=18):
    tf = textbox(s, left, top, width, height)
    first = True
    for item in items:
        lvl, text = (item if isinstance(item, tuple) else (0, item))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + text
        r.font.size = Pt(size - 3 * lvl)
        r.font.color.rgb = GREY if lvl else NAVY
        p.space_after = Pt(7)


def image(s, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        return
    kw = {}
    if width: kw["width"] = Inches(width)
    if height: kw["height"] = Inches(height)
    s.shapes.add_picture(path, Inches(left), Inches(top), **kw)


# ── 1. Title ────────────────────────────────────────────────────────────────
s = slide()
tf = textbox(s, 1.0, 2.2, 11.3, 2.6)
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Adaptive Packet Routing with Deep Reinforcement Learning"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = NAVY
p2 = tf.add_paragraph(); r2 = p2.add_run()
r2.text = "Topology generalization · multi-objective trade-offs · robustness · continual adaptation"
r2.font.size = Pt(19); r2.font.color.rgb = ACCENT
p3 = tf.add_paragraph(); p3.space_before = Pt(24)
r3 = p3.add_run(); r3.text = "Vaibhav Rathod"
r3.font.size = Pt(18); r3.font.color.rgb = GREY

# ── 2. The problem ──────────────────────────────────────────────────────────
s = slide(); title_bar(s, "The problem", "Why learn to route?")
bullets(s, [
    "Classical routing (OSPF/Dijkstra, ECMP) reacts to congestion and failures only via periodic metric updates — it does not learn.",
    "Reinforcement learning can learn a forwarding policy that optimizes long-horizon delay and delivery directly from link state.",
    "But the RL-routing literature is fragmented:",
    (1, "one algorithm, one topology, one objective, one aggregate number."),
    (1, "rarely answered together: generalization, trade-offs, robustness, continual adaptation, fairness."),
    "This project studies all of them on one simulator with one set of agents — honestly.",
], top=1.9, size=19)

# ── 3. Contributions ────────────────────────────────────────────────────────
s = slide(); title_bar(s, "Contributions")
bullets(s, [
    "Reproducible multi-seed comparison: 5 learners + 3 classical baselines, bootstrap CIs, 73-test suite.",
    "Zero-shot scalability: a GNN policy transfers route quality to graphs 5× larger than training.",
    "Multi-objective Pareto view — across algorithms and within one learner (reward-weight sweep).",
    "Adversarial robustness: worst-case (critical-link) vs. random failures.",
    "Continual learning: quantify catastrophic forgetting in routing + two remedies.",
    "All code, tests, and result artifacts released.",
], top=1.9, size=19)

# ── 4. System model ─────────────────────────────────────────────────────────
s = slide(); title_bar(s, "System model", "10-node network, realistic dynamics")
bullets(s, [
    "M/M/1 queuing-delay model: delay diverges as utilization → 1.",
    "AR(1) bursty congestion; RED-like loss; Markov link failure/recovery.",
    "Observation: 4 features/link (congestion, delay, loss, up/down) + src/dst.",
    "Agents: DQN, Rainbow, GNN-DQN, Q-Routing (+ continual variant).",
    "Baselines: Dijkstra (OSPF), ECMP, random.",
], left=0.7, top=1.9, width=6.2, size=18)
image(s, os.path.join(RESULTS, "topology.png"), 7.3, 1.9, width=5.6)

# ── 5. Headline result ──────────────────────────────────────────────────────
s = slide(); title_bar(s, "Headline: the GNN policy wins", "PDR & delay vs link-failure rate")
bullets(s, [
    "GNN-DQN: PDR = 1.00 at every failure rate, lowest delay (~9–10 ms).",
    "Routes from graph structure → reroutes around failures with no degradation.",
    "Rainbow competitive (0.90–0.94); vanilla DQN trails baselines on PDR.",
    "Caveat: failed links are penalized but traversable, so PDR saturates —",
    (1, "we add more discriminating metrics on the next slides."),
], left=0.7, top=1.9, width=6.0, size=17)
image(s, os.path.join(RESULTS, "evaluation_comparison.png"), 6.9, 1.7, width=6.1)

# ── 6. Scalability ──────────────────────────────────────────────────────────
s = slide(); title_bar(s, "Zero-shot scalability", "Trained on 10 nodes → applied to 50")
bullets(s, [
    "GNN weights are size-independent → apply to any graph.",
    "Metric: delay stretch = policy delay / Dijkstra-optimal.",
    "GNN stays ~1.1–1.2 up to 50 nodes (5× training size).",
    "Random walk degrades 3.2 → 10.7.",
    "Transfers route QUALITY, not just reachability.",
], left=0.7, top=1.9, width=5.6, size=18)
image(s, os.path.join(RESULTS, "scalability.png"), 6.5, 2.0, width=6.4)

# ── 7. Pareto ───────────────────────────────────────────────────────────────
s = slide(); title_bar(s, "Multi-objective trade-offs", "Latency vs reliability is not one number")
bullets(s, [
    "Across algorithms (delay, PDR): GNN-DQN is the SOLE Pareto-optimal point.",
    "Within one learner: sweep the reward's drop-penalty weight to trace its own front.",
    "Lets you pick an operating point instead of hard-coding one scalarization.",
], left=0.7, top=1.9, width=5.4, size=18)
image(s, os.path.join(RESULTS, "pareto.png"), 6.3, 2.0, width=6.6)

# ── 8. Adversarial ──────────────────────────────────────────────────────────
s = slide(); title_bar(s, "Adversarial robustness", "Worst-case vs average-case failures")
bullets(s, [
    "Targeted attack disables the highest edge-betweenness (most critical) links.",
    "GNN keeps PDR = 1.00 under both random and targeted failures.",
    "Vanilla DQN collapses 0.91 → 0.60 under attack —",
    (1, "an exposure invisible to random-failure testing."),
], left=0.7, top=1.9, width=5.6, size=18)
image(s, os.path.join(RESULTS, "adversarial.png"), 6.5, 2.0, width=6.4)

# ── 9. Continual learning (highlight) ───────────────────────────────────────
s = slide(); title_bar(s, "Catastrophic forgetting — the new result",
                       "Can an agent adapt without forgetting?")
bullets(s, [
    "Train one DQN sequentially on two conflicting destination tasks.",
    "Naive: Task-A delivery COLLAPSES 0.77 → 0.09 (catastrophic forgetting).",
    "Rehearsal (keep old replay data): fully prevents it (forgetting ≈ 0).",
    "EWC does NOT transfer here: destination-conditioned policy shares weights,",
    (1, "diagonal-Fisher can't isolate task-specific computation."),
    "Takeaway: data-space rehearsal (≈ free in RL) is the robust fix.",
], left=0.7, top=1.9, width=5.9, size=16)
image(s, os.path.join(RESULTS, "continual_learning.png"), 6.7, 2.0, width=6.3)

# ── 10. Honesty / limitations ───────────────────────────────────────────────
s = slide(); title_bar(s, "Honest limitations")
bullets(s, [
    "Failed links are penalized but traversable → PDR saturates; discriminating metrics carry the analysis.",
    "Most experiments use one 10-node topology (scalability covers size, GNN-only, no failures).",
    "Simulation only; a localhost data-plane prototype forwards real packets and reroutes, but no hardware testbed yet.",
    "We investigated a non-Markovian / recurrent-memory thesis, found it unsupported in this env, and removed it rather than overclaim.",
], top=1.9, size=18)

# ── 11. Wrap-up ─────────────────────────────────────────────────────────────
s = slide(); title_bar(s, "Takeaways & reproducibility")
bullets(s, [
    "A structure-aware GNN policy is Pareto-optimal, scales zero-shot, and resists targeted attacks.",
    "Catastrophic forgetting is real in learned routing; simple rehearsal is the robust remedy.",
    "Everything is reproducible: train.py, evaluate.py, experiments/*.py, 73 passing tests.",
    "Future work: larger & real testbeds, harder (hard-cut) failure regimes, decentralized multi-agent RL.",
], top=1.9, size=19)
tf = textbox(s, 0.7, 6.3, 12.0, 0.8)
r = tf.paragraphs[0].add_run()
r.text = "Code · tests · results · paper — all in the repository."
r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = ACCENT

prs.save(OUT)
print(f"Deck written -> {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
